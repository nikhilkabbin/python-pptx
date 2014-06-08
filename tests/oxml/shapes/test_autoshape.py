# encoding: utf-8

"""
Test suite for pptx.oxml.autoshape module.
"""

from __future__ import absolute_import, print_function

import pytest

from hamcrest import assert_that, instance_of, is_, none

from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsdecls
from pptx.oxml.shapes.autoshape import CT_PresetGeometry2D, CT_Shape
from pptx.oxml.shapes.shared import (
    CT_ShapeProperties, ST_Direction, ST_PlaceholderSize, ST_PlaceholderType
)

from ..unitdata.shape import (
    a_gd, a_prstGeom, an_avLst, an_spPr, test_shape_elements
)
from ...unitutil import TestCase


class DescribeCT_PresetGeometry2D(object):

    def it_can_get_the_gd_elms_as_a_sequence(self, gd_lst_fixture):
        prstGeom, expected_vals = gd_lst_fixture
        actual_vals = [(gd.name, gd.fmla) for gd in prstGeom.gd_lst]
        assert actual_vals == expected_vals

    def it_can_rewrite_the_gd_elms(self, rewrite_guides_fixture_):
        prstGeom, guides, expected_xml = rewrite_guides_fixture_
        print(prstGeom.xml)
        prstGeom.rewrite_guides(guides)
        assert prstGeom.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        [],
        [('adj1', 'val 111')],
        [('adj2', 'val 222'), ('adj3', 'val 333'), ('adj4', 'val 444'),
         ('adj5', 'val 555')],
    ])
    def gd_lst_fixture(self, request):
        expected_vals = request.param
        prstGeom = self.prstGeom_bldr('foobar', expected_vals).element
        return prstGeom, expected_vals

    @pytest.fixture(params=[
        ('circularArrow', 5),  # all five guides for five adj shape
        ('chevron',       0),  # empty guides for single adj shape
        ('chevron',       1),  # one guide for single adj shape
    ])
    def rewrite_guides_fixture_(self, request):
        prst, gd_count = request.param

        names = ('adj1', 'adj2', 'adj3', 'adj4', 'adj5')
        vals = (111, 222, 333, 444, 555)
        fmlas = [('val %d' % v) for v in vals]
        before_vals = [('adj6', 'val 666')] * 3
        after_vals = zip(names[:gd_count], fmlas[:gd_count])

        prstGeom = self.prstGeom_bldr(prst, before_vals).element
        guides = zip(names[:gd_count], vals[:gd_count])
        expected_xml = self.prstGeom_bldr(prst, after_vals).xml()

        return prstGeom, guides, expected_xml

    # fixture components ---------------------------------------------

    def prstGeom_bldr(self, prst, gd_vals):
        avLst_bldr = an_avLst()
        for name, fmla in gd_vals:
            gd_bldr = a_gd().with_name(name).with_fmla(fmla)
            avLst_bldr.with_child(gd_bldr)
        prstGeom_bldr = (
            a_prstGeom().with_nsdecls().with_prst(prst).with_child(
                avLst_bldr)
        )
        return prstGeom_bldr


class DescribeCT_Shape(object):

    def it_knows_its_MSO_AUTO_SHAPE_TYPE(self):
        rounded_rect_sp = test_shape_elements.rounded_rectangle
        placeholder_sp = test_shape_elements.placeholder
        # verify -----------------------
        assert rounded_rect_sp.prst == MSO_SHAPE.ROUNDED_RECTANGLE
        assert placeholder_sp.prst is None


class DescribeCT_ShapeProperties(object):

    def it_is_used_by_the_parser_for_an_spPr_element(self, spPr):
        assert isinstance(spPr, CT_ShapeProperties)

    # fixtures ---------------------------------------------

    @pytest.fixture
    def spPr(self):
        return an_spPr().with_nsdecls().element


class TestCT_Shape(TestCase):

    def test_is_autoshape_distinguishes_auto_shape(self):
        """CT_Shape.is_autoshape distinguishes auto shape"""
        # setup ------------------------
        autoshape = test_shape_elements.autoshape
        placeholder = test_shape_elements.placeholder
        textbox = test_shape_elements.textbox
        # verify -----------------------
        assert_that(autoshape.is_autoshape, is_(True))
        assert_that(placeholder.is_autoshape, is_(False))
        assert_that(textbox.is_autoshape, is_(False))

    def test_is_placeholder_distinguishes_placeholder(self):
        """CT_Shape.is_autoshape distinguishes placeholder"""
        # setup ------------------------
        autoshape = test_shape_elements.autoshape
        placeholder = test_shape_elements.placeholder
        textbox = test_shape_elements.textbox
        # verify -----------------------
        assert_that(autoshape.is_autoshape, is_(True))
        assert_that(placeholder.is_autoshape, is_(False))
        assert_that(textbox.is_autoshape, is_(False))

    def test_is_textbox_distinguishes_text_box(self):
        """CT_Shape.is_textbox distinguishes text box"""
        # setup ------------------------
        autoshape = test_shape_elements.autoshape
        placeholder = test_shape_elements.placeholder
        textbox = test_shape_elements.textbox
        # verify -----------------------
        assert_that(autoshape.is_textbox, is_(False))
        assert_that(placeholder.is_textbox, is_(False))
        assert_that(textbox.is_textbox, is_(True))

    def test_new_autoshape_sp_generates_correct_xml(self):
        """CT_Shape._new_autoshape_sp() returns correct XML"""
        # setup ------------------------
        id_ = 9
        name = 'Rounded Rectangle 8'
        prst = 'roundRect'
        left, top, width, height = 111, 222, 333, 444
        xml = (
            '<p:sp %s>\n  <p:nvSpPr>\n    <p:cNvPr id="%d" name="%s"/>\n    <'
            'p:cNvSpPr/>\n    <p:nvPr/>\n  </p:nvSpPr>\n  <p:spPr>\n    <a:xf'
            'rm>\n      <a:off x="%d" y="%d"/>\n      <a:ext cx="%d" cy="%d"/'
            '>\n    </a:xfrm>\n    <a:prstGeom prst="%s">\n      <a:avLst/>\n'
            '    </a:prstGeom>\n  </p:spPr>\n  <p:style>\n    <a:lnRef idx="1'
            '">\n      <a:schemeClr val="accent1"/>\n    </a:lnRef>\n    <a:f'
            'illRef idx="3">\n      <a:schemeClr val="accent1"/>\n    </a:fil'
            'lRef>\n    <a:effectRef idx="2">\n      <a:schemeClr val="accent'
            '1"/>\n    </a:effectRef>\n    <a:fontRef idx="minor">\n      <a:'
            'schemeClr val="lt1"/>\n    </a:fontRef>\n  </p:style>\n  <p:txBo'
            'dy>\n    <a:bodyPr rtlCol="0" anchor="ctr"/>\n    <a:lstStyle/>'
            '\n    <a:p>\n      <a:pPr algn="ctr"/>\n    </a:p>\n  </p:txBody'
            '>\n</p:sp>\n' %
            (nsdecls('a', 'p'), id_, name, left, top, width, height, prst)
        )
        # exercise ---------------------
        sp = CT_Shape.new_autoshape_sp(id_, name, prst, left, top,
                                       width, height)
        # verify -----------------------
        self.assertEqualLineByLine(xml, sp)

    def test_new_placeholder_sp_generates_correct_xml(self):
        """CT_Shape._new_placeholder_sp() returns correct XML"""
        # setup ------------------------
        expected_xml_tmpl = (
            '<p:sp %s>\n  <p:nvSpPr>\n    <p:cNvPr id="%s" name="%s"/>\n    <'
            'p:cNvSpPr>\n      <a:spLocks noGrp="1"/>\n    </p:cNvSpPr>\n    '
            '<p:nvPr>\n      <p:ph%s/>\n    </p:nvPr>\n  </p:nvSpPr>\n  <p:sp'
            'Pr/>\n%s</p:sp>\n' % (nsdecls('a', 'p'), '%d', '%s', '%s', '%s')
        )
        txBody_snippet = (
            '  <p:txBody>\n    <a:bodyPr/>\n    <a:lstStyle/>\n    <a:p/>\n  '
            '</p:txBody>\n')
        test_cases = (
            (2, 'Title 1', ST_PlaceholderType.CTR_TITLE, ST_Direction.HORZ,
             ST_PlaceholderSize.FULL, 0),
            (3, 'Date Placeholder 2', ST_PlaceholderType.DT,
             ST_Direction.HORZ, ST_PlaceholderSize.HALF, 10),
            (4, 'Vertical Subtitle 3', ST_PlaceholderType.SUB_TITLE,
             ST_Direction.VERT, ST_PlaceholderSize.FULL, 1),
            (5, 'Table Placeholder 4', ST_PlaceholderType.TBL,
             ST_Direction.HORZ, ST_PlaceholderSize.QUARTER, 14),
            (6, 'Slide Number Placeholder 5', ST_PlaceholderType.SLD_NUM,
             ST_Direction.HORZ, ST_PlaceholderSize.QUARTER, 12),
            (7, 'Footer Placeholder 6', ST_PlaceholderType.FTR,
             ST_Direction.HORZ, ST_PlaceholderSize.QUARTER, 11),
            (8, 'Content Placeholder 7', ST_PlaceholderType.OBJ,
             ST_Direction.HORZ, ST_PlaceholderSize.FULL, 15)
        )
        expected_values = (
            (2, 'Title 1', ' type="%s"' % ST_PlaceholderType.CTR_TITLE,
             txBody_snippet),
            (3, 'Date Placeholder 2', ' type="%s" sz="half" idx="10"' %
             ST_PlaceholderType.DT, ''),
            (4, 'Vertical Subtitle 3', ' type="%s" orient="vert" idx="1"' %
             ST_PlaceholderType.SUB_TITLE, txBody_snippet),
            (5, 'Table Placeholder 4', ' type="%s" sz="quarter" idx="14"' %
             ST_PlaceholderType.TBL, ''),
            (6, 'Slide Number Placeholder 5', ' type="%s" sz="quarter" '
             'idx="12"' % ST_PlaceholderType.SLD_NUM, ''),
            (7, 'Footer Placeholder 6', ' type="%s" sz="quarter" idx="11"' %
             ST_PlaceholderType.FTR, ''),
            (8, 'Content Placeholder 7', ' idx="15"', txBody_snippet)
        )
        # exercise ---------------------
        for case_idx, argv in enumerate(test_cases):
            id_, name, ph_type, orient, sz, idx = argv
            sp = CT_Shape.new_placeholder_sp(id_, name, ph_type, orient, sz,
                                             idx)
            # verify ------------------
            expected_xml = expected_xml_tmpl % expected_values[case_idx]
            self.assertEqualLineByLine(expected_xml, sp)

    def test_new_textbox_sp_generates_correct_xml(self):
        """CT_Shape.new_textbox_sp() returns correct XML"""
        # setup ------------------------
        id_ = 9
        name = 'TextBox 8'
        left, top, width, height = 111, 222, 333, 444
        xml = (
            '<p:sp %s>\n  <p:nvSpPr>\n    <p:cNvPr id="%d" name="%s"/>\n    <'
            'p:cNvSpPr txBox="1"/>\n    <p:nvPr/>\n  </p:nvSpPr>\n  <p:spPr>'
            '\n    <a:xfrm>\n      <a:off x="%d" y="%d"/>\n      <a:ext cx="%'
            'd" cy="%d"/>\n    </a:xfrm>\n    <a:prstGeom prst="rect">\n     '
            ' <a:avLst/>\n    </a:prstGeom>\n    <a:noFill/>\n  </p:spPr>\n  '
            '<p:txBody>\n    <a:bodyPr wrap="none">\n      <a:spAutoFit/>\n  '
            '  </a:bodyPr>\n    <a:lstStyle/>\n    <a:p/>\n  </p:txBody>\n</p'
            ':sp>\n' %
            (nsdecls('a', 'p'), id_, name, left, top, width, height)
        )
        # exercise ---------------------
        sp = CT_Shape.new_textbox_sp(id_, name, left, top, width, height)
        # verify -----------------------
        self.assertEqualLineByLine(xml, sp)

    def test_prstGeom_return_value(self):
        """CT_Shape.prstGeom value is correct"""
        # setup ------------------------
        rounded_rect_sp = test_shape_elements.rounded_rectangle
        placeholder_sp = test_shape_elements.placeholder
        # verify -----------------------
        assert_that(rounded_rect_sp.prstGeom,
                    instance_of(CT_PresetGeometry2D))
        assert_that(placeholder_sp.prstGeom, is_(none()))